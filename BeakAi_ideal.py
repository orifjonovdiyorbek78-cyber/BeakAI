# beak_ai_mobile_v2.py
# Beak AI — Mobile (Kivy) — Professional enhanced version
# Requirements: kivy, pillow, python-docx, python-pptx, openpyxl, fpdf, plyer
# Use Buildozer to make .apk (android). For iOS use kivy-ios pipeline.
import os
from pathlib import Path
from functools import partial
from io import BytesIO

from kivy.app import App
from kivy.lang import Builder
from kivy.uix.screenmanager import ScreenManager, Screen
from kivy.utils import get_color_from_hex
from kivy.uix.popup import Popup
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.label import Label
from kivy.uix.textinput import TextInput
from kivy.uix.button import Button
from kivy.metrics import dp
from kivy.core.window import Window
from kivy.clock import mainthread

# External libs
from PIL import Image, ImageDraw, ImageFont
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook, load_workbook

# For file chooser on Android / mobile
try:
    from plyer import filechooser
except Exception:
    filechooser = None

# -----------------------
# BASE_DIR for mobile storage
# -----------------------
def get_base_dir():
    try:
        # on Android, plyer.storagepath may exist, else fallback to common path
        from android.storage import primary_external_storage_path
        base = Path(primary_external_storage_path()) / "BeakAI_Office_Pro"
    except Exception:
        base = Path.home() / "BeakAI_Office_Pro"
    base.mkdir(parents=True, exist_ok=True)
    for s in ["Documents", "Presentations", "PDFs", "Excels", "Images", "Temp"]:
        (base / s).mkdir(exist_ok=True)
    return base

BASE_DIR = get_base_dir()

# -----------------------
# KV UI
# -----------------------
KV = r'''
#:import rgba kivy.utils.get_color_from_hex
<HeaderBar@BoxLayout>:
    size_hint_y: None
    height: dp(64)
    padding: dp(10)
    spacing: dp(8)
    canvas.before:
        Color:
            rgba: app.bg_color
        Rectangle:
            pos: self.pos
            size: self.size

    BoxLayout:
        orientation: 'vertical'
        size_hint_x: 0.85
        Label:
            text: "Beak AI"
            color: app.muted_color
            bold: True
            font_size: '18sp'
            size_hint_y: None
            height: self.texture_size[1]
        Label:
            text: "Created by Asadbek Tursunaliyev"
            color: app.muted_color
            font_size: '11sp'
            size_hint_y: None
            height: self.texture_size[1]

    Button:
        id: settings_btn
        text: "⚙"
        size_hint_x: 0.15
        on_release: app.open_settings()
        background_normal: ''
        background_color: app.small_btn_color
        color: app.btn_text_color

<MainScreen>:
    BoxLayout:
        orientation: 'vertical'
        canvas.before:
            Color:
                rgba: app.bg_color
            Rectangle:
                pos: self.pos
                size: self.size
        HeaderBar:
        BoxLayout:
            orientation: 'vertical'
            padding: dp(12)
            spacing: dp(10)

            Button:
                text: "🖼 Image → PDF (Ideal)"
                size_hint_y: None
                height: dp(64)
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
                on_release: app.open_image_pdf_picker()

            Button:
                text: "📄 Word Editor (.docx)"
                size_hint_y: None
                height: dp(64)
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
                on_release: root.manager.current = 'word'

            Button:
                text: "🎞 Presentation"
                size_hint_y: None
                height: dp(64)
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
                on_release: app.open_presentation_options()

            Button:
                text: "📊 Excel Editor (.xlsx)"
                size_hint_y: None
                height: dp(64)
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
                on_release: root.manager.current = 'excel'

        BoxLayout:
            size_hint_y: None
            height: dp(60)
            padding: dp(8)
            Button:
                text: "🤖 Chatbot"
                size_hint_x: 1
                background_normal: ''
                background_color: app.small_btn_color
                color: app.btn_text_color
                on_release: root.manager.current = 'chat'

        Label:
            text: "Beak AI - Bu kelajak sari ilk qadam"
            size_hint_y: None
            height: dp(36)
            color: app.muted_color
            font_size: '12sp'

<ImagePDFScreen>:
    BoxLayout:
        orientation: 'vertical'
        padding: dp(10)
        spacing: dp(10)
        Label:
            text: "Image → PDF (Tanlang va PDF yarating)"
            size_hint_y: None
            height: dp(40)
            color: app.muted_color
        ScrollView:
            id: scr
            GridLayout:
                id: file_grid
                cols: 1
                size_hint_y: None
                height: self.minimum_height
                row_default_height: dp(40)
        BoxLayout:
            size_hint_y: None
            height: dp(56)
            spacing: dp(8)
            Button:
                text: "Galeriyadan tanlash"
                on_release: root.open_gallery()
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
            Button:
                text: "PDF yaratish (1-100 tanlangan)"
                on_release: root.create_pdf()
                background_normal: ''
                background_color: app.small_btn_color
                color: app.btn_text_color
        Button:
            text: "Orqaga"
            size_hint_y: None
            height: dp(44)
            on_release: root.manager.current = 'main'
            background_normal: ''
            background_color: app.btn_color
            color: app.btn_text_color

<WordScreen>:
    BoxLayout:
        orientation: 'vertical'
        padding: dp(10)
        spacing: dp(8)
        Label:
            text: "Word Editor (.docx)"
            size_hint_y: None
            height: dp(36)
            color: app.muted_color
        TextInput:
            id: word_title
            hint_text: "Sarlavha (ixtiyoriy)"
            size_hint_y: None
            height: dp(40)
            font_name: 'Roboto'
            font_size: '16sp'
        ScrollView:
            TextInput:
                id: word_editor
                size_hint_y: None
                height: dp(300)
                text: ""
                font_name: 'Roboto'
                font_size: '14sp'
        BoxLayout:
            size_hint_y: None
            height: dp(56)
            spacing: dp(8)
            Button:
                text: "Rasm qo'shish"
                on_release: root.add_image_to_doc()
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
            Button:
                text: "Saqlash (.docx)"
                on_release: root.save_docx()
                background_normal: ''
                background_color: app.small_btn_color
                color: app.btn_text_color
        Button:
            text: "Orqaga"
            on_release: root.manager.current = 'main'
            size_hint_y: None
            height: dp(44)
            background_normal: ''
            background_color: app.btn_color
            color: app.btn_text_color

<PPTXEditorScreen>:
    BoxLayout:
        orientation: 'vertical'
        padding: dp(8)
        spacing: dp(6)
        Label:
            id: ppt_label
            text: root.header_text
            size_hint_y: None
            height: dp(36)
            color: app.muted_color
        BoxLayout:
            size_hint_y: None
            height: dp(40)
            spacing: dp(6)
            Label:
                text: "Slides:"
                size_hint_x: None
                width: dp(60)
                color: app.muted_color
            TextInput:
                id: slide_index
                text: "1"
                size_hint_x: None
                width: dp(60)
            Button:
                text: "Oldingi"
                size_hint_x: None
                width: dp(80)
                on_release: root.prev_slide()
            Button:
                text: "Keyingi"
                size_hint_x: None
                width: dp(80)
                on_release: root.next_slide()
        TextInput:
            id: slide_title
            hint_text: "Slide sarlavhasi"
            size_hint_y: None
            height: dp(40)
        ScrollView:
            TextInput:
                id: slide_text
                hint_text: "Slide matni"
                size_hint_y: None
                height: dp(180)
        BoxLayout:
            size_hint_y: None
            height: dp(48)
            spacing: dp(8)
            Button:
                text: "Rasm qo'shish"
                on_release: root.add_image_to_slide()
            Button:
                text: "Fon rangini o'zgartirish"
                on_release: root.change_bg_color()
            Button:
                text: "Shrift o'lchami"
                on_release: root.change_font_size()
        BoxLayout:
            size_hint_y: None
            height: dp(52)
            spacing: dp(8)
            Button:
                text: "Eksport PDF"
                on_release: root.export_pdf()
                background_normal: ''
                background_color: app.small_btn_color
                color: app.btn_text_color
            Button:
                text: "Eksport PPTX"
                on_release: root.export_pptx()
                background_normal: ''
                background_color: app.small_btn_color
                color: app.btn_text_color
        Button:
            text: "Orqaga"
            size_hint_y: None
            height: dp(44)
            on_release: root.manager.current = 'main'

<ExcelScreen>:
    BoxLayout:
        orientation: 'vertical'
        padding: dp(8)
        spacing: dp(8)
        Label:
            text: "Excel Editor (Simple)"
            size_hint_y: None
            height: dp(36)
            color: app.muted_color
        ScrollView:
            id: excel_scroll
            GridLayout:
                id: excel_grid
                cols: 10
                size_hint_y: None
                height: self.minimum_height
        BoxLayout:
            size_hint_y: None
            height: dp(44)
            spacing: dp(8)
            Button:
                text: "Saqlash (.xlsx)"
                on_release: root.save_xlsx()
                background_normal: ''
                background_color: app.small_btn_color
                color: app.btn_text_color
            Button:
                text: "Load (.xlsx)"
                on_release: root.load_xlsx()
                background_normal: ''
                background_color: app.btn_color
                color: app.btn_text_color
        Button:
            text: "Orqaga"
            size_hint_y: None
            height: dp(44)
            on_release: root.manager.current = 'main'

<ChatScreen>:
    BoxLayout:
        orientation: 'vertical'
        padding: dp(8)
        spacing: dp(6)
        Label:
            text: "Beak AI Chatbot"
            size_hint_y: None
            height: dp(36)
            color: app.muted_color
        ScrollView:
            GridLayout:
                id: chat_box
                cols: 1
                size_hint_y: None
                height: self.minimum_height
        BoxLayout:
            size_hint_y: None
            height: dp(48)
            spacing: dp(8)
            TextInput:
                id: user_msg
                multiline: False
                hint_text: "Xabar kiriting..."
            Button:
                text: "Yuborish"
                size_hint_x: None
                width: dp(100)
                on_release: root.send_msg(user_msg.text)
'''

# -----------------------
# Utility popup
# -----------------------
def popup(title, msg, width=0.9, height=None):
    content = BoxLayout(orientation='vertical', padding=10, spacing=10)
    content.add_widget(Label(text=msg))
    btn = Button(text='OK', size_hint_y=None, height=dp(44))
    content.add_widget(btn)
    p = Popup(title=title, content=content, size_hint=(width, None), height=dp(180) if not height else height)
    btn.bind(on_release=p.dismiss)
    p.open()

# -----------------------
# Screens Implementation
# -----------------------
class MainScreen(Screen):
    pass

class ImagePDFScreen(Screen):
    # holds selected files list in self.selected_files
    def on_enter(self):
        self.selected_files = []
        self.populate_grid()

    def populate_grid(self):
        grid = self.ids.file_grid
        grid.clear_widgets()
        if not hasattr(self, 'selected_files') or not self.selected_files:
            grid.add_widget(Label(text="Hali rasm tanlanmadi", size_hint_y=None, height=dp(40), color=App.get_running_app().muted_color))
        else:
            for p in self.selected_files:
                grid.add_widget(Label(text=str(Path(p).name), size_hint_y=None, height=dp(36), color=App.get_running_app().muted_color))

    def open_gallery(self):
        # Try plyer.filechooser (works on Android)
        if filechooser:
            try:
                # allow multiple
                filechooser.open_file(on_selection=self._on_files_selected, multiple=True, filters=['*.png','*.jpg','*.jpeg','*.webp','*.bmp'])
                return
            except Exception:
                pass
        # fallback to instruct user
        popup("Diqqat", "Galereyani ochib rasm tanlash uchun telefoningizdagi fayl tanlash oynasini ishlating yoki Buildozer orqali plyer filechooser qo'shing.")

    @mainthread
    def _on_files_selected(self, selection):
        if not selection:
            return
        # limit to 100
        sel = selection[:100]
        self.selected_files = sel
        self.populate_grid()

    def create_pdf(self):
        files = getattr(self, 'selected_files', [])
        if not files:
            popup("Diqqat", "Iltimos, kamida bitta rasm tanlang.")
            return
        if len(files) > 100:
            popup("Diqqat", "Maksimal 100 ta rasm qabul qilinadi.")
            files = files[:100]
        save_path = BASE_DIR / "PDFs" / f"images_to_pdf_{len(files)}.pdf"
        try:
            pil_images = []
            for p in files:
                img = Image.open(p).convert("RGB")
                max_w, max_h = 1240, 1754
                img.thumbnail((max_w, max_h), Image.LANCZOS)
                pil_images.append(img)
            pil_images[0].save(str(save_path), save_all=True, append_images=pil_images[1:], quality=95)
            popup("✅", f"PDF yaratildi:\n{save_path}")
        except Exception as e:
            popup("Xatolik", str(e))

class WordScreen(Screen):
    def on_enter(self):
        # ensure defaults
        self.last_added_image = None
    def add_image_to_doc(self):
        if filechooser:
            filechooser.open_file(on_selection=self._add_image_selected)
        else:
            popup("Diqqat", "Filechooser mavjud emas. Plyer filechooser ni buildozer spec ga qo‘shing.")

    @mainthread
    def _add_image_selected(self, selection):
        if not selection:
            return
        img_path = selection[0]
        # Insert marker into editor at end or store as list
        self.ids.word_editor.text += f"\n[IMAGE:{img_path}]\n"
        popup("✅", f"Rasm marker qo'shildi:\n{img_path}")

    def save_docx(self):
        title = self.ids.word_title.text.strip()
        content = self.ids.word_editor.text
        path = BASE_DIR / "Documents" / "beak_doc.docx"
        try:
            doc = Document()
            # set default style font
            style = doc.styles['Normal']
            style.font.name = 'Times New Roman'
            style.font.size = Pt(14)
            if title:
                doc.add_heading(title, level=1)
            lines = content.splitlines()
            for line in lines:
                line = line.strip()
                if line.startswith("[IMAGE:") and line.endswith("]"):
                    img_path = line[7:-1]
                    if os.path.exists(img_path):
                        try:
                            doc.add_picture(img_path, width=Inches(5))
                        except Exception:
                            doc.add_paragraph("[Rasm qo'shishda xatolik]")
                    else:
                        doc.add_paragraph("[Rasm topilmadi]")
                else:
                    p = doc.add_paragraph(line)
                    p.style = doc.styles['Normal']
            doc.save(str(path))
            popup("✅", f".docx saqlandi:\n{path}")
        except Exception as e:
            popup("Xatolik", str(e))

class PPTXEditorScreen(Screen):
    def on_enter(self):
        # default slides structure: list of dict {title,text,images:list,font_size:int,bg_color:hex}
        if not hasattr(self, 'slides'):
            self.slides = [{'title':'','text':'','images':[],'font_size':20,'bg_color':'#ffffff'}]
            self.current = 0
        self.update_ui()

    def header_text(self):
        return "Presentation Editor"
    header_text = property(lambda self: "Presentation Editor")

    def update_ui(self):
        s = self.slides[self.current]
        self.ids.slide_index.text = str(self.current+1)
        self.ids.slide_title.text = s.get('title','')
        self.ids.slide_text.text = s.get('text','')

    def prev_slide(self):
        if self.current > 0:
            self.current -= 1
            self.update_ui()

    def next_slide(self):
        if self.current < len(self.slides)-1:
            self.current += 1
            self.update_ui()

    def add_image_to_slide(self):
        if not filechooser:
            popup("Diqqat", "Filechooser mavjud emas. Plyer filechooser ni buildozer spec ga qo‘shing.")
            return
        filechooser.open_file(on_selection=self._on_image_selected)

    @mainthread
    def _on_image_selected(self, selection):
        if not selection:
            return
        img = selection[0]
        self.slides[self.current]['images'].append(img)
        popup("✅", f"Rasm qo'shildi:\n{img}")

    def change_bg_color(self):
        # simple color input popup
        def on_ok(instance):
            val = ti.text.strip()
            if val:
                self.slides[self.current]['bg_color'] = val
                popup("✅", "Fon rangi o'zgardi")
            popup_inst.dismiss()

        content = BoxLayout(orientation='vertical', spacing=6, padding=6)
        ti = TextInput(text=self.slides[self.current].get('bg_color','#ffffff'), hint_text="#rrggbb")
        btn = Button(text='OK', size_hint_y=None, height=dp(44))
        content.add_widget(Label(text="Hex rang kiriting (#rrggbb)"))
        content.add_widget(ti)
        content.add_widget(btn)
        popup_inst = Popup(title='Fon rangi', content=content, size_hint=(0.9,None), height=dp(220))
        btn.bind(on_release=on_ok)
        popup_inst.open()

    def change_font_size(self):
        def on_ok(instance):
            try:
                val = int(ti.text)
                if 8 <= val <= 72:
                    self.slides[self.current]['font_size'] = val
                    popup("✅", "Shrift o'lchami o'zgardi")
                else:
                    popup("Xatolik", "8 dan 72 gacha son kiriting")
            except:
                popup("Xatolik", "Butun son kiriting")
            popup_inst.dismiss()

        content = BoxLayout(orientation='vertical', spacing=6, padding=6)
        ti = TextInput(text=str(self.slides[self.current].get('font_size',20)), input_filter='int')
        btn = Button(text='OK', size_hint_y=None, height=dp(44))
        content.add_widget(Label(text="Shrift o'lchamini kiriting (8-72)"))
        content.add_widget(ti)
        content.add_widget(btn)
        popup_inst = Popup(title='Shrift o\'lchami', content=content, size_hint=(0.9,None), height=dp(220))
        btn.bind(on_release=on_ok)
        popup_inst.open()

    def export_pdf(self):
        # render slides as images and generate pdf
        try:
            images = []
            for i,s in enumerate(self.slides):
                # create PIL image for slide
                W,H = 1240, 700
                bg = s.get('bg_color','#ffffff')
                img = Image.new('RGB', (W,H), bg)
                draw = ImageDraw.Draw(img)
                # text
                title = s.get('title','')
                body = s.get('text','')
                # default font (may not match on mobile)
                try:
                    f_title = ImageFont.truetype("arial.ttf", 36)
                    f_body = ImageFont.truetype("arial.ttf", s.get('font_size',20))
                except:
                    f_title = ImageFont.load_default()
                    f_body = ImageFont.load_default()
                draw.text((40,40), title, font=f_title, fill=(0,0,0))
                draw.text((40,120), body, font=f_body, fill=(0,0,0))
                # paste images (first one if exists)
                imgs = s.get('images',[])
                y = 200
                for ipath in imgs[:3]:
                    try:
                        im = Image.open(ipath)
                        im.thumbnail((W-100, 300), Image.LANCZOS)
                        img.paste(im, (40,y))
                        y += im.size[1] + 20
                    except:
                        pass
                images.append(img)
            save_path = BASE_DIR / "PDFs" / "presentation_export.pdf"
            images[0].save(str(save_path), save_all=True, append_images=images[1:], quality=90)
            popup("✅", f"PDF eksport qilindi:\n{save_path}")
        except Exception as e:
            popup("Xatolik", str(e))

    def export_pptx(self):
        try:
            prs = Presentation()
            for s in self.slides:
                layout = prs.slide_layouts[6]  # blank
                slide = prs.slides.add_slide(layout)
                # bg color
                bg_color = s.get('bg_color','#FFFFFF').lstrip('#')
                try:
                    r = int(bg_color[0:2],16)/255.0
                    g = int(bg_color[2:4],16)/255.0
                    b = int(bg_color[4:6],16)/255.0
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = (int(r*255),int(g*255),int(b*255))
                except:
                    pass
                # add title
                if s.get('title'):
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = s.get('title','')
                    run.font.size = Pt(28)
                if s.get('text'):
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = s.get('text','')
                    run.font.size = Pt(s.get('font_size',20))
                # add images
                for ipath in s.get('images',[])[:3]:
                    try:
                        slide.shapes.add_picture(str(ipath), Inches(0.5), Inches(3.5), width=Inches(5))
                    except:
                        pass
            save_path = BASE_DIR / "Presentations" / "presentation_export.pptx"
            prs.save(str(save_path))
            popup("✅", f"PPTX eksport qilindi:\n{save_path}")
        except Exception as e:
            popup("Xatolik", str(e))

class ExcelScreen(Screen):
    def on_enter(self):
        # build a 10x10 grid of TextInput to simulate excel
        grid = self.ids.excel_grid
        grid.clear_widgets()
        self.cells = {}
        rows, cols = 10, 10
        # header row A..J
        for c in range(cols):
            lbl = Label(text=chr(65+c), size_hint_y=None, height=dp(36), color=App.get_running_app().muted_color)
            grid.add_widget(lbl)
        for r in range(rows):
            for c in range(cols):
                ti = TextInput(multiline=False, size_hint_y=None, height=dp(36))
                key = f"{chr(65+c)}{r+1}"
                self.cells[key] = ti
                grid.add_widget(ti)

    def save_xlsx(self):
        try:
            path = BASE_DIR / "Excels" / "beak_excel.xlsx"
            wb = Workbook()
            ws = wb.active
            for key, widget in self.cells.items():
                col = ord(key[0]) - 65 + 1
                row = int(key[1:])
                val = widget.text.strip()
                # simple formula handling: =SUM(A1:B2)
                if val.startswith('=SUM(') and val.endswith(')'):
                    rng = val[5:-1]
                    total = 0
                    try:
                        a,b = rng.split(':')
                        col1 = ord(a[0])-65
                        row1 = int(a[1:]) - 1
                        col2 = ord(b[0])-65
                        row2 = int(b[1:]) - 1
                        for rr in range(row1, row2+1):
                            for cc in range(col1, col2+1):
                                k = f"{chr(65+cc)}{rr+1}"
                                v = self.cells[k].text.strip()
                                try:
                                    total += float(v)
                                except:
                                    pass
                    except:
                        total = 0
                    ws.cell(row=row, column=col, value=total)
                else:
                    ws.cell(row=row, column=col, value=val)
            wb.save(str(path))
            popup("✅", f"Excel saqlandi:\n{path}")
        except Exception as e:
            popup("Xatolik", str(e))

    def load_xlsx(self):
        if not filechooser:
            popup("Diqqat", "Filechooser mavjud emas.")
            return
        filechooser.open_file(on_selection=self._on_xlsx_selected)

    @mainthread
    def _on_xlsx_selected(self, selection):
        if not selection:
            return
        p = selection[0]
        try:
            wb = load_workbook(p)
            ws = wb.active
            for r in ws.iter_rows(min_row=1, max_row=10, min_col=1, max_col=10, values_only=True):
                for c_idx, v in enumerate(r):
                    key = f"{chr(65+c_idx)}{r[0] and 1 or 1}"  # not used
            # simple: just reload into cells if sizes match
            for r_idx in range(1,11):
                for c_idx in range(1,11):
                    k = f"{chr(64+c_idx)}{r_idx}"
                    val = ws.cell(row=r_idx, column=c_idx).value
                    if k in self.cells:
                        self.cells[k].text = "" if val is None else str(val)
            popup("✅", "Excel yuklandi")
        except Exception as e:
            popup("Xatolik", str(e))

class ChatScreen(Screen):
    def send_msg(self, text):
        if not text or not text.strip():
            return
        self._write("Siz: " + text)
        reply = self.chat_bot_logic(text)
        self._write("🤖: " + reply)
        self.ids.user_msg.text = ""

    def _write(self, txt):
        l = Label(text=txt, size_hint_y=None, height=dp(28), color=App.get_running_app().muted_color)
        self.ids.chat_box.add_widget(l)

    def chat_bot_logic(self, msg):
        m = msg.lower()
        if "salom" in m:
            return "Salom! Qanday yordam kerak?"
        if "pdf" in m:
            return "Image→PDF bo'limiga o'ting va rasmlarni tanlang."
        if "word" in m:
            return "Word bo'limida rasm qo'shish va saqlash mumkin."
        return "Kechirasiz, men hozirgina offline yordamchiman. Men menyudan ishlashni maslahat beraman."

# -----------------------
# App
# -----------------------
class BeakAIApp(App):
    title = "Beak AI"
    # Colors
    bg_color = list(get_color_from_hex("#07102b"))  # dark navy
    btn_color = list(get_color_from_hex("#233b89"))
    small_btn_color = list(get_color_from_hex("#2a4aa0"))
    btn_text_color = (1, 1, 1, 1)
    muted_color = list(get_color_from_hex("#9fb1ff"))

    def build(self):
        self.sm = ScreenManager()
        Builder.load_string(KV)
        self.sm.add_widget(MainScreen(name='main'))
        self.sm.add_widget(ImagePDFScreen(name='image_pdf'))
        self.sm.add_widget(WordScreen(name='word'))
        self.sm.add_widget(PPTXEditorScreen(name='pptx_editor'))
        self.sm.add_widget(PPTXEditorScreen(name='pptx'))
        self.sm.add_widget(ExcelScreen(name='excel'))
        self.sm.add_widget(ChatScreen(name='chat'))
        return self.sm

    def on_start(self):
        pass

    def open_image_pdf_picker(self):
        # go to screen and auto open gallery
        self.sm.current = 'image_pdf'
        screen = self.sm.get_screen('image_pdf')
        # if plyer available open gallery directly
        if filechooser:
            try:
                filechooser.open_file(on_selection=screen._on_files_selected, multiple=True, filters=['*.png','*.jpg','*.jpeg','*.bmp','*.webp'])
                return
            except Exception:
                pass
        # else just notify
        popup("Diqqat", "Rasm tanlash uchun Galereyani oching (plyer.filechooser ishlamadi).")

    def open_presentation_options(self):
        # popup with two options: PDF or PPTX and ask for slide count (1-50)
        def on_pdf(instance):
            popup_inst.dismiss()
            self._ask_slide_count(export_type='pdf')

        def on_pptx(instance):
            popup_inst.dismiss()
            self._ask_slide_count(export_type='pptx')

        content = BoxLayout(orientation='vertical', spacing=8, padding=8)
        b1 = Button(text="PDF taqdimot yaratish", size_hint_y=None, height=dp(48))
        b2 = Button(text="PPTX taqdimot yaratish", size_hint_y=None, height=dp(48))
        content.add_widget(b1); content.add_widget(b2)
        popup_inst = Popup(title="Taqdimot turi", content=content, size_hint=(0.9,None), height=dp(180))
        b1.bind(on_release=on_pdf); b2.bind(on_release=on_pptx)
        popup_inst.open()

    def _ask_slide_count(self, export_type='pdf'):
        def on_ok(inst):
            try:
                cnt = int(ti.text)
                if not (1 <= cnt <= 50):
                    raise ValueError
            except:
                popup("Xatolik", "1 dan 50 gacha son kiriting")
                return
            popup2.dismiss()
            # create editor screen with cnt slides
            screen = self.sm.get_screen('pptx')
            screen.slides = [{'title':'','text':'','images':[],'font_size':20,'bg_color':'#ffffff'} for _ in range(cnt)]
            screen.current = 0
            screen.update_ui()
            self.sm.current = 'pptx'
        content = BoxLayout(orientation='vertical', spacing=8, padding=8)
        ti = TextInput(text='4', input_filter='int', multiline=False)
        btn = Button(text='OK', size_hint_y=None, height=dp(44))
        content.add_widget(Label(text=f"Slidelar soni ({export_type.upper()}) — 1 dan 50 gacha:"))
        content.add_widget(ti)
        content.add_widget(btn)
        popup2 = Popup(title="Slide sonini kiriting", content=content, size_hint=(0.9,None), height=dp(220))
        btn.bind(on_release=on_ok)
        popup2.open()

    def open_settings(self):
        # settings popup: theme and language
        def toggle_theme(inst):
            # toggle simple dark/light by swapping colors
            if self.bg_color[0] == 0.03:  # dark approx
                self.bg_color = get_color_from_hex("#f4f7ff") + [1]
                self.btn_text_color = (0,0,0,1)
            else:
                self.bg_color = get_color_from_hex("#07102b") + [1]
                self.btn_text_color = (1,1,1,1)
            popup3.dismiss()

        def set_lang(inst):
            popup3.dismiss()
            popup("Til o'zgardi", "Til muvaffaqiyatli o'zgartirildi (test).")

        content = BoxLayout(orientation='vertical', spacing=8, padding=8)
        btn_theme = Button(text='Yorug\'/Qorong\'u rejimni almashtirish', size_hint_y=None, height=dp(48))
        btn_lang = Button(text='Til (UZ/EN/RU)', size_hint_y=None, height=dp(48))
        content.add_widget(btn_theme); content.add_widget(btn_lang)
        popup3 = Popup(title="Sozlamalar", content=content, size_hint=(0.9,None), height=dp(220))
        btn_theme.bind(on_release=toggle_theme); btn_lang.bind(on_release=set_lang)
        popup3.open()

# -----------------------
# Run
# -----------------------
if __name__ == '__main__':
    BeakAIApp().run()